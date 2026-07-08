#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Dung file main.pptx (PowerPoint sua duoc) tu noi dung bo slide S3 (4 Act).
Chay:  python3 build_pptx.py   (trong thu muc main/)
Yeu cau: pip install python-pptx ; file figure3.png cung thu muc.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))

# ---- Palette ----
ACCENT  = RGBColor(0x25, 0x63, 0xEB)
ACCENT2 = RGBColor(0x08, 0x91, 0xB2)
GOOD    = RGBColor(0x16, 0xA3, 0x4A)
WARN    = RGBColor(0xD9, 0x77, 0x06)
BAD     = RGBColor(0xDC, 0x26, 0x26)
INK     = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG_IDEA = RGBColor(0xEF, 0xF6, 0xFF)
BG_TECH = RGBColor(0xEC, 0xFE, 0xFF)
BG_KEY  = RGBColor(0xF0, 0xFD, 0xF4)
BG_WARN = RGBColor(0xFE, 0xF2, 0xF2)
BG_GREY = RGBColor(0xF8, 0xFA, 0xFC)

FONT = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
    return s

def _set_text(tf, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """blocks: list of paragraphs; each paragraph = list of runs (text,size,bold,color)."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        for (text, size, bold, color) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT

def textbox(s, l, t, w, h, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set_text(tb.text_frame, blocks, align, anchor)
    return tb

def box(s, l, t, w, h, fill, bar, label, body_runs, label_size=12, body_size=14):
    """Rounded rect voi thanh mau trai + label in dam + body."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = bar; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    blocks = []
    if label:
        blocks.append([(label, label_size, True, bar)])
    blocks.extend(body_runs)
    _set_text(tf, blocks)
    return sh

def heading(s, text, act=None):
    textbox(s, 0.55, 0.32, 9.6, 0.9, [[(text, 30, True, ACCENT)]])
    if act:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(0.42), Inches(2.15), Inches(0.5))
        b.fill.solid(); b.fill.fore_color.rgb = ACCENT2; b.line.fill.background(); b.shadow.inherit = False
        _set_text(b.text_frame, [[(act, 12, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def section(title, subtitle, tag):
    s = slide(ACCENT)
    textbox(s, 1.0, 2.5, 11.3, 1.2, [[(tag, 16, True, RGBColor(0xBF,0xDB,0xFE))]], align=PP_ALIGN.CENTER)
    textbox(s, 1.0, 3.1, 11.3, 1.6, [[(title, 44, True, WHITE)]], align=PP_ALIGN.CENTER)
    textbox(s, 1.5, 4.7, 10.3, 1.2, [[(subtitle, 20, False, RGBColor(0xE0,0xF2,0xFE))]], align=PP_ALIGN.CENTER)
    return s

def bullets(s, l, t, w, h, items, size=16, color=INK):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb

def table(s, l, t, w, rows, col_w=None, header_fill=ACCENT, fs=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(0.4*nr)).table
    if col_w:
        for j, cw in enumerate(col_w): gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.margin_left = Inches(0.08); c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.name = FONT
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                r.font.bold = True; r.font.color.rgb = WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else BG_GREY
                r.font.color.rgb = INK
    return gt

# ============================================================
# S1 — Title
# ============================================================
s = slide(BG_IDEA)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SW, Inches(0.09))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, 1.0, 1.15, 11.3, 0.6, [[("ACL 2025 · Long Papers", 15, True, ACCENT2)]], align=PP_ALIGN.CENTER)
textbox(s, 1.0, 1.6, 11.3, 1.1, [[("S³ — Semantic Signal Separation", 40, True, ACCENT)]], align=PP_ALIGN.CENTER)
textbox(s, 1.0, 2.95, 11.3, 0.7, [[("Phân tích chủ đề của văn bản sử dụng thuật toán ICA", 22, False, MUTED)]], align=PP_ALIGN.CENTER)
textbox(s, 1.5, 3.9, 10.3, 0.9,
        [[("Làm sao tách được từng ", 16, False, INK), ("giọng nói", 16, True, INK),
          (" trong một căn phòng ồn ào? — và điều đó liên quan gì đến việc tìm ", 16, False, INK),
          ("chủ đề", 16, True, INK), (" trong văn bản?", 16, False, INK)]], align=PP_ALIGN.CENTER)
textbox(s, 1.0, 5.3, 11.3, 0.5, [[("Kardos, Kostkan, Enevoldsen, Vermillet, Nielbo, Rocca · Aarhus University", 13, False, MUTED)]], align=PP_ALIGN.CENTER)

# ============================================================
# ACT 1
# ============================================================
section("Tách tín hiệu", "Cocktail party → văn bản → ICA vs PCA", "ACT 1 · HOOK")

# S2 cocktail
s = slide(); heading(s, 'Bài toán "bữa tiệc cocktail"', "ACT 1")
textbox(s, 0.55, 1.35, 12.2, 0.9, [[("Ba người nói cùng lúc; mỗi micro chỉ thu được một ", 16, False, INK),
    ("mớ âm thanh trộn lẫn", 16, True, INK), (". Thuật toán ", 16, False, INK), ("ICA", 16, True, ACCENT),
    (" tách lại từng giọng — dù không biết trước giọng ai (blind source separation).", 16, False, INK)]])
# flow
flow = [("Giọng A / B / C", BG_IDEA, ACCENT), ("TRỘN\n(micro)", BG_WARN, BAD), ("ICA\ngỡ trộn", BG_KEY, GOOD), ("→ A / B / C", BG_KEY, GOOD)]
x = 0.9
for i,(txt,fill,bar) in enumerate(flow):
    b = box(s, x, 2.6, 2.5, 1.3, fill, bar, "", [[(txt, 15, True, INK)]])
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in b.text_frame.paragraphs: p.alignment = PP_ALIGN.CENTER
    x += 3.0
    if i < len(flow)-1:
        textbox(s, x-0.55, 2.95, 0.5, 0.6, [[("→", 28, True, MUTED)]], align=PP_ALIGN.CENTER)
box(s, 0.9, 4.35, 11.5, 1.0, BG_TECH, ACCENT2, "3 Ý CẦN NHỚ",
    [[("(1) nhiều nguồn ", 14, False, INK), ("độc lập", 14, True, INK), ("; (2) cái ta quan sát là ", 14, False, INK),
      ("bản trộn", 14, True, INK), ("; (3) vẫn ", 14, False, INK), ("gỡ ngược", 14, True, INK), (" ra từng nguồn.", 14, False, INK)]])

# S3 bridge
s = slide(); heading(s, 'Văn bản cũng là một "bản trộn"', "ACT 1")
box(s, 0.55, 1.5, 6.0, 3.2, BG_IDEA, ACCENT, "ĐỔI GÓC NHÌN",
    [[("• Một tài liệu hiếm khi nói đúng một chủ đề — nó trộn nhiều chủ đề (như phòng tiệc trộn nhiều giọng).", 15, False, INK)],
     [("• Tự động tìm & tách các chủ đề trong kho văn bản = topic modeling.", 15, False, INK)]])
box(s, 6.8, 1.5, 5.95, 3.2, BG_KEY, GOOD, "Ý TƯỞNG CỦA S³",
    [[("Coi mỗi chủ đề như một \"giọng nói\" → dùng đúng ICA để tách chủ đề ra khỏi văn bản.", 15, False, INK)],
     [("Đó là lý do tên bài: Semantic Signal Separation — \"tách tín hiệu ngữ nghĩa\".", 15, True, INK)]])
box(s, 0.55, 5.0, 12.2, 1.2, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Đổi \"giọng nói\" thành \"chủ đề\", \"căn phòng\" thành \"một bài viết\". Nếu tách được chủ đề như tách giọng nói thì sao? Đó đúng là điều S³ làm.", 14, False, INK)]])

# S4 ICA vs PCA
s = slide(); heading(s, "ICA khác PCA thế nào? (ý chính)", "ACT 1")
box(s, 0.55, 1.5, 6.0, 2.1, BG_TECH, ACCENT2, "PCA",
    [[("Tìm hướng dữ liệu trải rộng nhất; các trục chỉ cần không tương quan — điều kiện dễ. Hay dùng để nén / giảm chiều.", 15, False, INK)]])
box(s, 6.8, 1.5, 5.95, 2.1, BG_KEY, GOOD, "ICA",
    [[("Tìm các thành phần độc lập thống kê — điều kiện khó hơn. Nhờ đó tách được các nguồn thật sự tách bạch.", 15, False, INK)]])
box(s, 0.55, 3.9, 12.2, 1.5, BG_IDEA, ACCENT, "CHỐT MỘT CÂU",
    [[("\"Không tương quan\" chưa chắc \"độc lập\". Chính vì đòi độc lập mà ICA cho các chủ đề ít chồng lấn → gốc rễ khiến S³ đạt diversity gần như hoàn hảo.", 15, False, INK)]])
textbox(s, 0.55, 5.6, 12.2, 0.5, [[("(Chi tiết y=x², whitening, công thức — để dành Act 3.)", 12, False, MUTED)]], align=PP_ALIGN.CENTER)

# ============================================================
# ACT 2
# ============================================================
section("Động lực & Bối cảnh", "Vì sao cần S³ — và so với ai", "ACT 2 · ĐỘNG LỰC")

# S5 topic model
s = slide(); heading(s, "Topic model là gì?", "ACT 2")
box(s, 0.55, 1.5, 12.2, 1.3, BG_IDEA, ACCENT, "ĐỊNH NGHĨA",
    [[("Nhóm phương pháp thống kê, unsupervised (không cần nhãn) tự khám phá chủ đề ẩn trong kho văn bản lớn. Mỗi topic = một tập từ khoá đại diện.", 15, False, INK)]])
box(s, 0.55, 3.0, 6.0, 1.4, BG_KEY, GOOD, "Topic thể thao", [[("bóng đá · cầu thủ · bàn thắng · HLV", 14, False, INK)]])
box(s, 6.8, 3.0, 5.95, 1.4, BG_KEY, GOOD, "Topic công nghệ", [[("AI · chip · dữ liệu · phần mềm", 14, False, INK)]])
box(s, 0.55, 4.7, 12.2, 1.2, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Có 100.000 bài báo, muốn biết chúng nói về gì mà không đọc hết — topic model trả về vài chục nhóm từ khoá đại diện.", 14, False, INK)]])

# S6 BoW
s = slide(); heading(s, "Cách cổ điển: đếm từ (Bag-of-Words)", "ACT 2")
textbox(s, 0.55, 1.35, 12.2, 0.6, [[("LSA (1988), LDA (2003) đều dựa trên bag-of-words: đếm tần suất từ, bỏ thứ tự & ngữ cảnh.", 15, False, INK)]])
box(s, 0.55, 2.1, 12.2, 2.5, BG_WARN, BAD, "3 GIỚI HẠN CỦA BoW (kẻ thù S³ muốn đánh bại)",
    [[("1. Nhạy với stop words (\"the, of, và, là\") → topic vô nghĩa nếu không làm sạch.", 15, False, INK)],
     [("2. Preprocessing tạo nhiều \"bậc tự do cho nhà nghiên cứu\" → khó tái lập.", 15, False, INK)],
     [("3. BoW thưa & số chiều cao → tính kém, khớp kém.", 15, False, INK)]])
box(s, 0.55, 4.9, 12.2, 1.0, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Đếm từ thì \"the, is, of\" luôn nhiều nhất, nên buộc phải làm sạch rất kỹ — và mỗi người làm một kiểu.", 14, False, INK)]])

# S7 embeddings
s = slide(); heading(s, "Embeddings mở cơ hội — nhưng vẫn vướng", "ACT 2")
box(s, 0.55, 1.5, 6.0, 3.0, BG_KEY, GOOD, "EMBEDDING GIÚP GÌ",
    [[("• Có ngữ cảnh → không cần vứt bỏ thông tin", 14, False, INK)],
     [("• Bền lỗi chính tả & từ lạ", 14, False, INK)],
     [("• Không gian liên tục → giả định Gaussian, cho transfer learning", 14, False, INK)]])
box(s, 6.8, 1.5, 5.95, 3.0, BG_WARN, BAD, "CONTEXTUAL MODEL HIỆN TẠI VẪN",
    [[("• Chậm, không ổn định", 14, False, INK)],
     [("• Vẫn cần preprocessing nặng", 14, False, INK)],
     [("• Không rõ có thật sự dùng ngữ cảnh — vì vẫn test trên dữ liệu đã làm sạch!", 14, False, INK)]])
box(s, 0.55, 4.8, 12.2, 1.1, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Ba vấn đề bên phải chính là thứ S³ giải quyết.", 14, True, INK)]])

# S8 baselines
s = slide(); heading(s, "Bối cảnh: các baseline để so sánh", "ACT 2")
table(s, 0.55, 1.5, 12.2, [
    ["Nhóm", "Mô hình", "Cơ chế cốt lõi", "Contextual?"],
    ["Classical", "LDA / NMF", "Xác suất / phân rã trên BoW", "Không"],
    ["Neural/VAE", "ZeroShotTM, CombinedTM", "VAE trên embedding (+BoW)", "Có"],
    ["Neural", "ECRTM / FASTopic", "Clustering reg. / Optimal transport", "Không / Có"],
    ["Clustering", "Top2Vec / BERTopic", "UMAP + HDBSCAN + (cosine / c-TF-IDF)", "Có"],
    ["Đề xuất", "S³", "ICA trên embedding", "Có"],
], col_w=[1.8, 3.0, 5.4, 2.0], fs=12)
box(s, 0.55, 5.3, 12.2, 1.0, BG_IDEA, ACCENT, "CHỈ CẦN NHỚ",
    [[("3 trường phái cũ: đếm từ, mạng nơ-ron, gom cụm. S³ đi hướng thứ tư — phân rã ma trận trên embedding.", 14, False, INK)]])

# S9 topic = truc
s = slide(); heading(s, "Cú lật: topic = trục", "ACT 2")
box(s, 0.55, 1.5, 3.9, 1.2, BG_GREY, MUTED, "", [[("BERTopic/Top2Vec → topic = cụm", 14, False, INK)]])
box(s, 4.65, 1.5, 3.9, 1.2, BG_GREY, MUTED, "", [[("LDA/CTM → topic = phân phối xác suất", 14, False, INK)]])
box(s, 8.75, 1.5, 4.0, 1.2, BG_KEY, ACCENT, "", [[("S³ → topic = TRỤC ngữ nghĩa", 14, True, ACCENT)]])
box(s, 0.55, 3.0, 12.2, 2.5, BG_IDEA, ACCENT, "TRỰC GIÁC",
    [[("Mỗi topic là một hướng giải thích phần biến thiên riêng của kho. Phân rã thành:", 15, False, INK)],
     [("• A = các topic (latent components)", 15, False, INK)],
     [("• S = độ mạnh mỗi topic trong từng tài liệu", 15, False, INK)],
     [("Dùng ICA để các trục độc lập.", 15, True, INK)]])

# ============================================================
# ACT 3
# ============================================================
section("Thuật toán S³", "6 bước · whitening · word importance · từ negative", "ACT 3 · THUẬT TOÁN")

# S10 6 steps
s = slide(); heading(s, "Thuật toán S³ — 6 bước", "ACT 3")
box(s, 0.55, 1.4, 12.2, 0.95, BG_TECH, ACCENT2, "CỐT LÕI",
    [[("Toàn bộ phần \"khó\" chỉ là một phép phân rã ma trận:  tài liệu = topic × độ mạnh", 16, True, INK)]])
steps = ["1. Encode tài liệu → ma trận X",
         "2. FastICA: X = A·S  (A = topic, S = document-topic importance)",
         "3. Encode từ vựng bằng cùng encoder → V",
         "4. Unmixing: C = A⁺",
         "5. Chiếu từ: W = V·Cᵀ",
         "6. Tính word importance cho từng topic"]
y = 2.6
for st in steps:
    box(s, 0.9, y, 11.5, 0.62, BG_GREY, RGBColor(0xE2,0xE8,0xF0), "", [[(st, 14, False, INK)]])
    y += 0.72

# S11 whitening
s = slide(); heading(s, "Chi tiết 1–2: phân rã + whitening", "ACT 3")
textbox(s, 0.55, 1.35, 12.2, 0.6, [[("X = A·S  bằng FastICA.", 20, True, INK)]], align=PP_ALIGN.CENTER)
for i,(lab,txt,fill,bar) in enumerate([
    ("①", "Gốc: elip nghiêng", BG_IDEA, ACCENT),
    ("②", "Whitened: cầu, var=1", BG_KEY, GOOD),
    ("③", "ICA xoay về trục độc lập", BG_TECH, ACCENT2)]):
    b = box(s, 0.9 + i*4.05, 2.2, 3.7, 1.3, fill, bar, lab, [[(txt, 14, True, INK)]])
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
box(s, 0.55, 3.9, 12.2, 2.0, BG_WARN, BAD, "PCA và ICA phối hợp (trả bài Act 1)",
    [[("FastICA là mô hình noiseless → phải whitening trước (dùng PCA). Ngay đó S³ giảm chiều còn N (giữ N principal components đầu). Rồi ICA mới xoay cho độc lập.", 15, False, INK)],
     [("→ PCA dọn dẹp, ICA tách. S³ KHÔNG chỉ là PCA.", 15, True, INK)]])

# S12 project
s = slide(); heading(s, "Chi tiết 3–5: chiếu từ lên trục", "ACT 3")
box(s, 0.9, 1.5, 11.5, 0.62, BG_GREY, RGBColor(0xE2,0xE8,0xF0), "", [[("3. Encode từ vựng → V", 15, False, INK)]])
box(s, 0.9, 2.22, 11.5, 0.62, BG_GREY, RGBColor(0xE2,0xE8,0xF0), "", [[("4. Unmixing matrix: C = A⁺  (pseudo-inverse của A)", 15, False, INK)]])
box(s, 0.9, 2.94, 11.5, 0.62, BG_GREY, RGBColor(0xE2,0xE8,0xF0), "", [[("5. Chiếu từ: W = V·Cᵀ", 15, False, INK)]])
box(s, 0.55, 3.85, 12.2, 1.0, BG_IDEA, ACCENT, "Ý NGHĨA",
    [[("Wⱼₜ = vị trí của từ j trên trục topic t → \"toạ độ\" để chọn từ mô tả topic.", 15, False, INK)]])
box(s, 0.55, 5.0, 12.2, 1.0, BG_TECH, ACCENT2, "INFERENCE tài liệu mới",
    [[("Chỉ một phép nhân ma trận: Ŝ = X̂·Cᵀ — cực nhanh.", 15, False, INK)]])

# S13 word importance + figure3
s = slide(); heading(s, "Bước 6: word importance (Figure 3)", "ACT 3")
img = os.path.join(HERE, "figure3.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(0.7), Inches(1.5), height=Inches(3.0))
box(s, 8.1, 1.5, 4.65, 3.0, BG_TECH, ACCENT2, "ĐỌC HÌNH",
    [[("vₜ = hình chiếu của từ lên trục topic", 13, False, INK)],
     [("‖v‖ = độ dài vector từ", 13, False, INK)],
     [("Θ = góc giữa v và trục", 13, False, INK)]])
box(s, 0.7, 4.75, 3.9, 1.5, BG_TECH, ACCENT2, "① AXIAL",
    [[("βₜⱼ = Wⱼₜ", 15, True, INK)], [("Từ nổi bật nhất. Coherence cao nhất.", 12, False, MUTED)]])
box(s, 4.75, 4.75, 3.9, 1.5, BG_TECH, ACCENT2, "② ANGULAR",
    [[("βₜⱼ = Wⱼₜ / ‖Wⱼ‖", 15, True, INK)], [("Từ đặc trưng nhất. Diversity cao nhất.", 12, False, MUTED)]])
box(s, 8.8, 4.75, 3.95, 1.5, BG_KEY, GOOD, "③ COMBINED (mặc định)",
    [[("βₜⱼ = (Wⱼₜ)³ / ‖Wⱼ‖", 15, True, INK)], [("Mũ lẻ để giữ dấu. Cân bằng ①②.", 12, False, MUTED)]])

# S14 negative
s = slide(); heading(s, "Điểm độc đáo: mô tả bằng từ NEGATIVE", "ACT 3")
box(s, 0.55, 1.5, 12.2, 1.5, BG_IDEA, ACCENT, "TRỰC GIÁC",
    [[("Cả 3 công thức cho phép từ có độ quan trọng âm → ngoài \"top từ dương\", S³ liệt kê cả top từ âm = một định nghĩa âm của topic.", 15, False, INK)]])
bullets(s, 0.7, 3.2, 12.0, 2.0, [
    "LSA cũng có tính chất này về toán, nhưng chưa từng được khai thác.",
    "Lợi ích: hai topic có từ-dương giống nhau vẫn phân biệt được nhờ từ âm.",
    "So sánh định lượng bỏ qua từ âm (để công bằng); từ âm chỉ dùng minh hoạ định tính.",
], size=15)
box(s, 0.55, 5.3, 12.2, 1.0, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Không mô hình nào khác cho biết một topic KHÔNG nói về cái gì — S³ thì có.", 14, False, INK)]])

# ============================================================
# ACT 4
# ============================================================
section("Thí nghiệm & Kết quả", "Setup · con số · hồi quy · preprocessing · định tính", "ACT 4 · CHỐT HẠ")

# S15 setup
s = slide(); heading(s, "Thiết lập thí nghiệm", "ACT 4")
table(s, 0.55, 1.4, 7.0, [
    ["Dataset", "#Tài liệu", "Vocab"],
    ["ArXiv ML / BBC", "2.048 / 1.225", "2.849 / 3.851"],
    ["20NG Preprocessed", "16.310", "1.612"],
    ["20NG Raw (thô)", "18.846", "21.668"],
    ["StackExchange / Wiki Med", "75.000 / 6.861", "17.884 / 22.145"],
], col_w=[3.4, 2.0, 1.6], fs=11)
box(s, 7.8, 1.4, 4.95, 2.5, BG_IDEA, ACCENT, "THIẾT KẾ TINH TẾ",
    [[("Dùng cả 20NG thô lẫn đã làm sạch để đo riêng ảnh hưởng preprocessing. 4 embeddings (GloVe→E5), topic 10–50, không tune hyperparameter, chạy 1 lần.", 13, False, INK)]])
box(s, 0.55, 4.3, 12.2, 1.7, BG_KEY, GOOD, "ĐO CHẤT LƯỢNG",
    [[("Diversity (d): topic khác nhau đến đâu · Coherence (C): topic mạch lạc đến đâu.", 14, False, INK)],
     [("C̄ = √(C_ex · C_in) ;  Interpretability = √(C̄ · d)  — dùng geometric mean để phạt lệch, buộc tốt cả hai mặt.", 14, False, INK)]])

# S16 results
s = slide(); heading(s, "Kết quả: các con số biết nói", "ACT 4")
for i,(num,cap) in enumerate([("4.5×","nhanh hơn á quân BERTopic"),("27.5×","nhanh hơn baselines (TB)"),("#1","hiệu năng tổng hợp")]):
    textbox(s, 0.7+i*4.1, 1.7, 3.8, 1.2, [[(num, 54, True, ACCENT)]], align=PP_ALIGN.CENTER)
    textbox(s, 0.7+i*4.1, 2.95, 3.8, 0.7, [[(cap, 13, False, MUTED)]], align=PP_ALIGN.CENTER)
box(s, 0.55, 4.0, 12.2, 2.2, BG_KEY, GOOD, "CÂN BẰNG TỐI ƯU",
    [[("• ECRTM, FASTopic: đa dạng hơn nhưng kém mạch lạc", 15, False, INK)],
     [("• Top2Vec: rất mạch lạc nhưng kém đa dạng", 15, False, INK)],
     [("• S³: cân bằng tối ưu coherence ↔ diversity", 15, True, GOOD)]])

# S17 regression
s = slide(); heading(s, "Kiểm định thống kê (rất mạnh)", "ACT 4")
box(s, 0.55, 1.4, 12.2, 0.85, BG_KEY, GOOD, "",
    [[("F = 167.4 · p < 0.001 · R² = 0.673  → S³ vượt trội có ý nghĩa thống kê so với mọi baseline.", 15, True, INK)]])
table(s, 0.55, 2.5, 12.2, [
    ["Mô hình", "Hệ số", "Mô hình", "Hệ số"],
    ["S³_com (mốc)", "0.606", "ZeroShotTM", "−0.117"],
    ["ECRTM", "−0.031", "CombinedTM", "−0.120"],
    ["FASTopic", "−0.043", "BERTopic", "−0.214"],
    ["Top2Vec", "−0.046", "NMF / LDA", "−0.222 / −0.272"],
], col_w=[3.5, 2.6, 3.5, 2.6], fs=12)
textbox(s, 0.55, 5.6, 12.2, 0.5, [[("Càng âm = càng kém S³. 3 biến thể S³ đứng đầu; LDA/NMF/BERTopic ở đáy.", 12, False, MUTED)]])

# S18 preprocessing
s = slide(); heading(s, "Phát hiện phản trực giác nhất", "ACT 4")
box(s, 0.55, 1.5, 12.2, 1.2, BG_KEY, GOOD, "",
    [[("S³ là mô hình DUY NHẤT liên tục chạy TỐT HƠN trên văn bản thô so với dữ liệu đã tiền xử lý.", 18, True, INK)]])
bullets(s, 0.7, 3.0, 12.0, 2.0, [
    "Baseline nhìn chung ngang hoặc kém đi khi bỏ preprocessing.",
    "S³ hưởng lợi nhiều nhất → thật sự khai thác thông tin thêm trong văn bản thô.",
    "Trên kho thô, S³ cao hơn TẤT CẢ — kể cả mô hình train trên dữ liệu đã làm sạch.",
], size=15)
box(s, 0.55, 5.3, 12.2, 1.1, RGBColor(0xFE,0xFC,0xE9), WARN, "CÂU DẪN",
    [[("Mọi mô hình khác cần bạn làm sạch dữ liệu; S³ ngược lại — càng để nguyên càng tốt. Bằng chứng nó thật sự dùng ngữ cảnh.", 14, False, INK)]])

# S19 qualitative
s = slide(); heading(s, "Định tính + Concept Compass", "ACT 4")
box(s, 0.55, 1.45, 12.2, 0.9, BG_WARN, BAD, "KÉM — LDA, BERTopic",
    [[("that · to · you · of · the · and   (toàn function words)", 13, False, INK)]])
box(s, 0.55, 2.45, 12.2, 0.9, BG_KEY, GOOD, "TỐT NHẤT — Top2Vec & S³",
    [[("epilepsy · medical · toxins · homeopathy   (S³ — y khoa)", 13, False, INK)]])
box(s, 0.55, 3.5, 12.2, 2.6, BG_IDEA, ACCENT, "CONCEPT COMPASS",
    [[("Đặt 2 trục lên mặt phẳng 2D → \"bản đồ ngữ nghĩa\" để đặt bất kỳ từ nào lên và xem nó nằm đâu giữa các trục.", 15, False, INK)],
     [("Topic 0 & 4 có từ dương giống nhau (đều \"clustering\") — chỉ từ ÂM mới phân biệt được.", 15, False, INK)],
     [("→ S³ không chỉ liệt kê topic — nó cho một bản đồ ngữ nghĩa. Thứ mà mô hình cụm/xác suất không làm được.", 15, True, INK)]])

# S20 conclusion
s = slide(); heading(s, "Kết luận & Hạn chế", "ACT 4")
box(s, 0.55, 1.5, 3.9, 1.6, BG_KEY, GOOD, "Chất lượng", [[("Mạch lạc + đa dạng, sạch", 14, False, INK)]])
box(s, 4.65, 1.5, 3.9, 1.6, BG_KEY, GOOD, "Tốc độ", [[("Nhanh nhất trong contextual TM", 14, False, INK)]])
box(s, 8.75, 1.5, 4.0, 1.6, BG_KEY, GOOD, "Không preprocessing", [[("Chạy tốt hơn trên text thô", 14, False, INK)]])
box(s, 0.55, 3.4, 12.2, 1.6, BG_TECH, ACCENT2, "HẠN CHẾ (khách quan)",
    [[("Metric có giả định mạnh · baseline được cài lại (CTM lệch nhỏ) · không tune hyperparameter · chỉ 1 seed · preprocessing thử trên 1 kho.", 14, False, INK)]])
box(s, 0.55, 5.2, 12.2, 1.0, BG_IDEA, ACCENT, "ĐỌNG LẠI",
    [[("Phương pháp đơn giản, theory-driven + thư viện Turftopic (MIT).", 15, False, INK)]])

# S21 Q&A
s = slide(ACCENT)
textbox(s, 1.0, 2.6, 11.3, 1.4, [[("Cảm ơn & Q&A", 48, True, WHITE)]], align=PP_ALIGN.CENTER)
textbox(s, 1.0, 4.1, 11.3, 0.8, [[("pip install turftopic  →  SemanticSignalSeparation(n_components=10)", 18, False, RGBColor(0xE0,0xF2,0xFE))]], align=PP_ALIGN.CENTER)
textbox(s, 1.0, 5.2, 11.3, 0.6, [[("S³ — Semantic Signal Separation · ACL 2025 · Aarhus University", 14, False, RGBColor(0xBF,0xDB,0xFE))]], align=PP_ALIGN.CENTER)

out = os.path.join(HERE, "main.pptx")
prs.save(out)
print("Da tao:", out, "-", len(prs.slides._sldIdLst), "slides")
